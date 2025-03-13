from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# Crear una nueva presentación
prs = Presentation()

# --- SLIDE 1: PORTADA ---
slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Layout "Título" (Title Slide)
slide1.shapes.title.text = "Predicción de Chollos en Portátiles"
subtitle = slide1.placeholders[1]
fecha = datetime.now().strftime("%d/%m/%Y")
subtitle.text = f"Análisis predictivo mediante web scraping y machine learning\nFecha: {fecha}\n[Logos del proyecto/instituto]"

# --- SLIDE 2: OBJETIVO DEL PROYECTO ---
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Layout "Título y contenido"
slide2.shapes.title.text = "Objetivo del Proyecto"
tf2 = slide2.shapes.placeholders[1].text_frame
tf2.text = "Desarrollar un sistema automático que identifique 'chollos' en tiendas online"
p = tf2.add_paragraph()
p.text = "Ayudar a los consumidores a encontrar las mejores relaciones calidad-precio"
p = tf2.add_paragraph()
p.text = "Automatizar la extracción, procesamiento y análisis de datos de portátiles"

# --- SLIDE 3: OBTENCIÓN DE DATOS - WEB SCRAPING ---
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Obtención de Datos - Web Scraping"
tf3 = slide3.shapes.placeholders[1].text_frame
tf3.text = "Fuentes de datos: PCBox y AppInformática"
p = tf3.add_paragraph()
p.text = "Herramientas utilizadas: Python, BeautifulSoup, Selenium"
p = tf3.add_paragraph()
p.text = "Proceso:"
# Sub-bullets para el proceso
p = tf3.add_paragraph()
p.text = "Identificación de URLs y estructura de las páginas"
p.level = 1
p = tf3.add_paragraph()
p.text = "Extracción de características: procesador, RAM, almacenamiento, precio, etc."
p.level = 1
p = tf3.add_paragraph()
p.text = "Gestión de paginación y categorías de productos"
p.level = 1
p = tf3.add_paragraph()
p.text = "Incluir capturas del código y resultados del scraping"
p.level = 1

# --- SLIDE 4: ALMACENAMIENTO DE DATOS ---
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Almacenamiento de Datos"
tf4 = slide4.shapes.placeholders[1].text_frame
tf4.text = "Diseño de la base de datos: esquema relacional"
p = tf4.add_paragraph()
p.text = "Tecnología: MySQL/PostgreSQL"
p = tf4.add_paragraph()
p.text = "Proceso ETL:"
p = tf4.add_paragraph()
p.text = "• Transformación de datos extraídos a formato estructurado"
p.level = 1
p = tf4.add_paragraph()
p.text = "• Limpieza inicial y detección de duplicados"
p.level = 1
p = tf4.add_paragraph()
p.text = "• Carga en la base de datos mediante script automatizado"
p.level = 1
p = tf4.add_paragraph()
p.text = "Incluir diagrama de la base de datos"
p.level = 1

# --- SLIDE 5: PREPARACIÓN DEL DATASET ---
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "Preparación del Dataset"
tf5 = slide5.shapes.placeholders[1].text_frame
tf5.text = "Consultas SQL para construir el dataset inicial"
p = tf5.add_paragraph()
p.text = "Unificación de datos de diferentes fuentes"
p = tf5.add_paragraph()
p.text = "Normalización y estandarización de formatos"
p = tf5.add_paragraph()
p.text = "Generación de variable objetivo 'Chollo' basada en criterios predefinidos"
p = tf5.add_paragraph()
p.text = "Mostrar la estructura final del dataset"

# --- SLIDE 6: EXPLORACIÓN Y LIMPIEZA DE DATOS ---
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
slide6.shapes.title.text = "Exploración y Limpieza de Datos"
tf6 = slide6.shapes.placeholders[1].text_frame
tf6.text = "Estadísticas descriptivas y visualización de distribuciones clave"
p = tf6.add_paragraph()
p.text = "Tratamiento de valores nulos:"
p.level = 0
p = tf6.add_paragraph()
p.text = "• Imputación de valores de batería basada en características similares"
p.level = 1
p = tf6.add_paragraph()
p.text = "• Gestión de campos vacíos en especificaciones técnicas"
p.level = 1
p = tf6.add_paragraph()
p.text = "Detección y tratamiento de outliers"
p.level = 0
p = tf6.add_paragraph()
p.text = "Incluir gráficos relevantes del análisis exploratorio"
p.level = 0

# --- SLIDE 7: INGENIERÍA DE CARACTERÍSTICAS ---
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
slide7.shapes.title.text = "Ingeniería de Características"
tf7 = slide7.shapes.placeholders[1].text_frame
tf7.text = "Creación de nuevas variables:"
p = tf7.add_paragraph()
p.text = "• Ratio precio/especificaciones"
p.level = 1
p = tf7.add_paragraph()
p.text = "• Variables dummy para sistemas operativos y tipos de procesador"
p.level = 1
p = tf7.add_paragraph()
p.text = "• Normalización de resoluciones de pantalla"
p.level = 1
p = tf7.add_paragraph()
p.text = "Codificación de variables categóricas y selección de características mediante análisis de correlación"
p.level = 0
p = tf7.add_paragraph()
p.text = "Incluir heatmap u otros gráficos relevantes"
p.level = 0

# --- SLIDE 8: MODELOS DE MACHINE LEARNING ---
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
slide8.shapes.title.text = "Modelos de Machine Learning"
tf8 = slide8.shapes.placeholders[1].text_frame
tf8.text = "Modelos implementados:"
p = tf8.add_paragraph()
p.text = "• Random Forest, Gradient Boosting, SVM, Regresión Logística"
p = tf8.add_paragraph()
p.text = "Evaluación mediante validación cruzada y comparación de métricas"
p = tf8.add_paragraph()
p.text = "Incluir gráficos de comparación entre modelos"

# --- SLIDE 9: RESULTADOS Y MODELO FINAL ---
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
slide9.shapes.title.text = "Resultados y Modelo Final"
tf9 = slide9.shapes.placeholders[1].text_frame
tf9.text = "Modelo seleccionado y justificación"
p = tf9.add_paragraph()
p.text = "Hiperparámetros optimizados"
p = tf9.add_paragraph()
p.text = "Métricas en conjunto de prueba: Accuracy, Precision, Recall, F1-score"
p = tf9.add_paragraph()
p.text = "Matriz de confusión y características más importantes"
p = tf9.add_paragraph()
p.text = "Incluir visualizaciones de la matriz de confusión e importancia de features"

# --- SLIDE 10: CONCLUSIONES Y TRABAJO FUTURO ---
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
slide10.shapes.title.text = "Conclusiones y Trabajo Futuro"
tf10 = slide10.shapes.placeholders[1].text_frame
tf10.text = "Resumen de logros y limitaciones encontradas"
p = tf10.add_paragraph()
p.text = "Posibles mejoras:"
p.level = 0
p = tf10.add_paragraph()
p.text = "• Ampliar fuentes de datos"
p.level = 1
p = tf10.add_paragraph()
p.text = "• Implementar técnicas avanzadas de deep learning"
p.level = 1
p = tf10.add_paragraph()
p.text = "• Desarrollo de una API para consultas en tiempo real"
p.level = 1
p = tf10.add_paragraph()
p.text = "Agradecimientos y contacto"
p.level = 0

# Guardar la presentación en un archivo
prs.save("presentacion_chollos.pptx")
print("Presentación creada y guardada como 'presentacion_chollos.pptx'")
